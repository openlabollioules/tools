"""
Analyse des templates de slides PowerPoint
Permet d'identifier les shapes, indices et placeholders pour chaque layout
"""

import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pathlib import Path

class SlideTemplateAnalyzer:
    def __init__(self):
        self.base_template_path = "./templates/"
        self.templates = {
            "french": {
                "public": "fr/CS-PU-template_fr.pptx",
                "internal": "fr/CS-IN-template_fr.pptx", 
                "confidential": "fr/CS-CO-template_fr.pptx"
            },
            "english": {
                "public": "en/CS-PU-template_en.pptx",
                "internal": "en/CS-IN-template_en.pptx",
                "confidential": "en/CS-CO-template_en.pptx"
            }
        }
        
    def get_shape_type_name(self, shape_type):
        """Convertit le type de shape en nom lisible"""
        try:
            # Essayer d'abord les types les plus communs
            common_types = {
                1: "AUTO_SHAPE",
                2: "CALLOUT", 
                3: "CHART",
                4: "COMMENT",
                5: "CONNECTOR",
                7: "EMBEDDED_OLE_OBJECT",
                8: "FORM_CONTROL",
                9: "FREEFORM",
                10: "GROUP",
                11: "IGX_GRAPHIC",
                12: "LINKED_OLE_OBJECT",
                13: "LINKED_PICTURE",
                14: "MEDIA",
                15: "OLE_CONTROL_OBJECT",
                16: "PICTURE",
                17: "PLACEHOLDER",
                18: "SCRIPT_ANCHOR",
                19: "SHAPE_TYPE_MIXED",
                20: "TABLE",
                21: "TEXT_EFFECT",
                22: "TEXT_BOX",
                23: "WEB_VIDEO"
            }
            
            # Si le type est un entier, utiliser le mapping direct
            if isinstance(shape_type, int):
                return common_types.get(shape_type, f"UNKNOWN_TYPE_{shape_type}")
            
            # Sinon, essayer d'utiliser l'énumération
            shape_types = {}
            
            # Ajouter seulement les types qui existent dans cette version de python-pptx
            try:
                shape_types[MSO_SHAPE_TYPE.AUTO_SHAPE] = "AUTO_SHAPE"
            except AttributeError:
                pass
            try:
                shape_types[MSO_SHAPE_TYPE.CALLOUT] = "CALLOUT"
            except AttributeError:
                pass
            try:
                shape_types[MSO_SHAPE_TYPE.CHART] = "CHART"
            except AttributeError:
                pass
            try:
                shape_types[MSO_SHAPE_TYPE.COMMENT] = "COMMENT"
            except AttributeError:
                pass
            try:
                shape_types[MSO_SHAPE_TYPE.EMBEDDED_OLE_OBJECT] = "EMBEDDED_OLE_OBJECT"
            except AttributeError:
                pass
            try:
                shape_types[MSO_SHAPE_TYPE.FORM_CONTROL] = "FORM_CONTROL"
            except AttributeError:
                pass
            try:
                shape_types[MSO_SHAPE_TYPE.FREEFORM] = "FREEFORM"
            except AttributeError:
                pass
            try:
                shape_types[MSO_SHAPE_TYPE.GROUP] = "GROUP"
            except AttributeError:
                pass
            try:
                shape_types[MSO_SHAPE_TYPE.PICTURE] = "PICTURE"
            except AttributeError:
                pass
            try:
                shape_types[MSO_SHAPE_TYPE.PLACEHOLDER] = "PLACEHOLDER"
            except AttributeError:
                pass
            try:
                shape_types[MSO_SHAPE_TYPE.TABLE] = "TABLE"
            except AttributeError:
                pass
            try:
                shape_types[MSO_SHAPE_TYPE.TEXT_BOX] = "TEXT_BOX"
            except AttributeError:
                pass
                
            return shape_types.get(shape_type, f"UNKNOWN_ENUM_{shape_type}")
            
        except Exception as e:
            return f"ERROR_GETTING_TYPE_{shape_type}"
    
    def get_placeholder_type_name(self, placeholder_type):
        """Convertit le type de placeholder en nom lisible"""
        from pptx.enum.shapes import PP_PLACEHOLDER
        placeholder_types = {
            PP_PLACEHOLDER.BODY: "BODY",
            PP_PLACEHOLDER.CHART: "CHART", 
            PP_PLACEHOLDER.CLIP_ART: "CLIP_ART",
            PP_PLACEHOLDER.CENTER_TITLE: "CENTER_TITLE",
            PP_PLACEHOLDER.CONTENT: "CONTENT",
            PP_PLACEHOLDER.DATE: "DATE",
            PP_PLACEHOLDER.FOOTER: "FOOTER",
            PP_PLACEHOLDER.HEADER: "HEADER",
            PP_PLACEHOLDER.MEDIA_CLIP: "MEDIA_CLIP",
            PP_PLACEHOLDER.OBJECT: "OBJECT",
            PP_PLACEHOLDER.ORG_CHART: "ORG_CHART",
            PP_PLACEHOLDER.PICTURE: "PICTURE",
            PP_PLACEHOLDER.SLIDE_NUMBER: "SLIDE_NUMBER",
            PP_PLACEHOLDER.SUBTITLE: "SUBTITLE",
            PP_PLACEHOLDER.TABLE: "TABLE",
            PP_PLACEHOLDER.TITLE: "TITLE",
            PP_PLACEHOLDER.VERTICAL_BODY: "VERTICAL_BODY",
            PP_PLACEHOLDER.VERTICAL_OBJECT: "VERTICAL_OBJECT",
            PP_PLACEHOLDER.VERTICAL_TITLE: "VERTICAL_TITLE"
        }
        return placeholder_types.get(placeholder_type, f"UNKNOWN_{placeholder_type}")
    
    def analyze_shape(self, shape, shape_index):
        """Analyse un shape spécifique"""
        try:
            shape_info = {
                "index": shape_index,
                "name": getattr(shape, 'name', 'No name'),
                "type": "UNKNOWN",
                "has_text_frame": hasattr(shape, 'text_frame') and shape.text_frame is not None,
                "is_placeholder": shape.is_placeholder if hasattr(shape, 'is_placeholder') else False,
                "width": 'N/A',
                "height": 'N/A',
                "left": 'N/A',
                "top": 'N/A'
            }
            
            # Essayer d'obtenir le type de shape
            try:
                shape_info["type"] = self.get_shape_type_name(shape.shape_type)
            except Exception as e:
                shape_info["type"] = f"ERROR_TYPE_{e}"
            
            # Essayer d'obtenir les dimensions
            try:
                if hasattr(shape, 'width'):
                    shape_info["width"] = shape.width.inches
                if hasattr(shape, 'height'):
                    shape_info["height"] = shape.height.inches
                if hasattr(shape, 'left'):
                    shape_info["left"] = shape.left.inches
                if hasattr(shape, 'top'):
                    shape_info["top"] = shape.top.inches
            except Exception as e:
                print(f"[DEBUG] Erreur dimensions pour shape {shape_index}: {e}")
            
            # Informations sur le placeholder
            if shape_info["is_placeholder"]:
                try:
                    placeholder = shape.placeholder_format
                    shape_info["placeholder_type"] = self.get_placeholder_type_name(placeholder.type)
                    shape_info["placeholder_idx"] = placeholder.idx
                except Exception as e:
                    shape_info["placeholder_type"] = f"ERROR_PLACEHOLDER_{e}"
                    shape_info["placeholder_idx"] = "N/A"
            
            # Informations sur le texte
            if shape_info["has_text_frame"]:
                try:
                    text_frame = shape.text_frame
                    shape_info["text_content"] = text_frame.text if text_frame.text else "Empty"
                    shape_info["paragraphs_count"] = len(text_frame.paragraphs)
                except Exception as e:
                    shape_info["text_content"] = f"ERROR_TEXT_{e}"
                    shape_info["paragraphs_count"] = 0
                
            return shape_info
            
        except Exception as e:
            return {
                "index": shape_index,
                "name": "ERROR_ANALYZING_SHAPE",
                "type": f"ERROR_{e}",
                "error": str(e),
                "has_text_frame": False,
                "is_placeholder": False,
                "width": 'N/A',
                "height": 'N/A',
                "left": 'N/A',
                "top": 'N/A'
            }
    
    def analyze_slide_layout(self, layout, layout_index):
        """Analyse un layout de slide"""
        layout_info = {
            "index": layout_index,
            "name": layout.name,
            "shapes_count": len(layout.shapes),
            "placeholders_count": len(layout.placeholders),
            "shapes": [],
            "placeholders": []
        }
        
        # Analyse des shapes
        for i, shape in enumerate(layout.shapes):
            shape_info = self.analyze_shape(shape, i)
            layout_info["shapes"].append(shape_info)
        
        # Analyse des placeholders
        for i, placeholder in enumerate(layout.placeholders):
            placeholder_info = self.analyze_shape(placeholder, i)
            placeholder_info["placeholder_index"] = i
            layout_info["placeholders"].append(placeholder_info)
            
        return layout_info
    
    def analyze_template(self, template_path):
        """Analyse un template PowerPoint complet"""
        if not os.path.exists(template_path):
            return {"error": f"Template not found: {template_path}"}
            
        try:
            prs = Presentation(template_path)
            template_info = {
                "path": template_path,
                "layouts_count": len(prs.slide_layouts),
                "layouts": []
            }
            
            # Analyse de chaque layout
            for i, layout in enumerate(prs.slide_layouts):
                layout_info = self.analyze_slide_layout(layout, i)
                template_info["layouts"].append(layout_info)
                
            return template_info
            
        except Exception as e:
            return {"error": f"Error analyzing template {template_path}: {str(e)}"}
    
    def print_analysis(self, template_info):
        """Affiche l'analyse de manière formatée"""
        if "error" in template_info:
            print(f"❌ {template_info['error']}")
            return
            
        print(f"\n📊 ANALYSE DU TEMPLATE: {template_info['path']}")
        print(f"📋 Nombre de layouts: {template_info['layouts_count']}")
        print("=" * 80)
        
        for layout in template_info["layouts"]:
            print(f"\n🎯 Layout {layout['index']}: {layout['name']}")
            print(f"   📦 Shapes: {layout['shapes_count']} | 🎭 Placeholders: {layout['placeholders_count']}")
            
            # Affichage des placeholders
            if layout["placeholders"]:
                print("\n   🎭 PLACEHOLDERS:")
                for ph in layout["placeholders"]:
                    print(f"      [{ph['placeholder_index']}] {ph['placeholder_type']} (idx: {ph.get('placeholder_idx', 'N/A')})")
                    if ph["has_text_frame"]:
                        print(f"          📝 Text: '{ph['text_content'][:50]}...' ({ph['paragraphs_count']} paragraphs)")
                    print(f"          📐 Position: ({ph['left']}, {ph['top']}) | Size: {ph['width']}x{ph['height']}")
            
            # Affichage des shapes non-placeholder
            non_placeholder_shapes = [s for s in layout["shapes"] if not s["is_placeholder"]]
            if non_placeholder_shapes:
                print("\n   📦 OTHER SHAPES:")
                for shape in non_placeholder_shapes:
                    print(f"      [{shape['index']}] {shape['name']} ({shape['type']})")
                    if shape["has_text_frame"]:
                        print(f"          📝 Text: '{shape['text_content'][:50]}...'")
                    print(f"          📐 Position: ({shape['left']}, {shape['top']}) | Size: {shape['width']}x{shape['height']}")
            
            print("-" * 60)
    
    def generate_code_suggestions(self, template_info):
        """Génère des suggestions de code basées sur l'analyse"""
        if "error" in template_info:
            return
            
        print(f"\n💡 SUGGESTIONS DE CODE pour {template_info['path']}:")
        print("=" * 80)
        
        for layout in template_info["layouts"]:
            print(f"\n# Layout {layout['index']}: {layout['name']}")
            print(f"# Utilisation: prs.slide_layouts[{layout['index']}]")
            
            if layout["placeholders"]:
                print("\n# Placeholders disponibles:")
                for ph in layout["placeholders"]:
                    if ph["placeholder_type"] == "TITLE":
                        print(f"# slide.shapes[{ph['index']}].text = 'Votre titre'  # {ph['placeholder_type']}")
                    elif ph["placeholder_type"] in ["BODY", "CONTENT"]:
                        print(f"# slide.shapes.placeholders[{ph['placeholder_index']}].text_frame  # {ph['placeholder_type']}")
                    elif ph["placeholder_type"] == "SUBTITLE":
                        print(f"# slide.shapes[{ph['index']}].text = 'Votre sous-titre'  # {ph['placeholder_type']}")
            
            print()
    
    def analyze_all_templates(self):
        """Analyse tous les templates disponibles"""
        print("🔍 ANALYSE DE TOUS LES TEMPLATES")
        print("=" * 80)
        
        for language, templates in self.templates.items():
            print(f"\n🌐 LANGUE: {language.upper()}")
            
            for confidentiality, template_file in templates.items():
                full_path = os.path.join(self.base_template_path, template_file)
                print(f"\n🔒 Confidentialité: {confidentiality.upper()}")
                
                template_info = self.analyze_template(full_path)
                self.print_analysis(template_info)
                self.generate_code_suggestions(template_info)

def test_single_template():
    """Fonction pour tester un seul template (pour debug)"""
    analyzer = SlideTemplateAnalyzer()
    
    # Tester avec un template spécifique - remplacez le chemin selon vos besoins
    test_paths = [
        "./templates/fr/CS-PU-template_fr.pptx",
        "./templates/en/CS-PU-template_en.pptx"
    ]
    
    for test_path in test_paths:
        if os.path.exists(test_path):
            print(f"\n🧪 TEST D'UN SEUL TEMPLATE: {test_path}")
            template_info = analyzer.analyze_template(test_path)
            analyzer.print_analysis(template_info)
            analyzer.generate_code_suggestions(template_info)
            break
    else:
        print("❌ Aucun template de test trouvé. Vérifiez les chemins.")

def main():
    """Fonction principale pour lancer l'analyse"""
    analyzer = SlideTemplateAnalyzer()
    
    print("🚀 Démarrage de l'analyse des templates PowerPoint...")
    
    # Vérifier si le dossier templates existe
    if not os.path.exists(analyzer.base_template_path):
        print(f"❌ Le dossier templates n'existe pas: {analyzer.base_template_path}")
        print("📁 Veuillez créer le dossier et y placer vos fichiers .pptx")
        return
    
    # Option 1: Analyser tous les templates
    # analyzer.analyze_all_templates()
    
    # Option 2: Tester un seul template (pour debug)
    test_single_template()
    
    print("\n✅ Analyse terminée!")
    print("\n💡 Utilisez les suggestions de code ci-dessus pour implémenter vos méthodes.")

if __name__ == "__main__":
    main() 