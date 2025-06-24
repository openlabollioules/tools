"""
title: Generate PPTX Presentation
author: openlab
version: 0.1
license: MIT
description: Génère un fichier PPTX via un LLM (Ollama) et renvoie un lien de téléchargement
"""

import os, uuid, tempfile
from typing import Optional
from pathlib import Path
from fastapi import UploadFile
import re
import json
from datetime import datetime
from pptx import Presentation
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches

from open_webui.routers.files import upload_file
from open_webui.models.users import Users
from fastapi import Request
from open_webui.storage.provider import Storage
from open_webui.models.files import Files, FileForm
from typing import Callable, Any

class EventEmitter:
    def __init__(self, event_emitter: Callable[[dict], Any] = None):
        self.event_emitter = event_emitter

    async def emit(self, description="Unknown State", status="in_progress", done=False):
        if self.event_emitter:
            await self.event_emitter(
                {
                    "type": "status",
                    "data": {
                        "status": status,
                        "description": description,
                        "done": done,
                    },
                }
            )

class HelpFunctions:
    def __init__(self):
        
        self.slide_layouts = {
            "title_and_content": 1,
            "abstract": 2,
            "chapter_title": 3,
            "basic_content": 4,
            "final_slide_fr": 12,
            "final_slide_en": 13,
        }

    def remove_tags_no_keep(self, text : str, start : str, end : str) -> str:
        """
        Remove all text between two tags (`start` and `end`), tags included.

        Parameters
        ----------
        text : str
            text to remove tags from
        start : str
            starting tag
        end : str
            ending tag

        Returns
        -------
        str
            Text with tags removed
        """
        return re.sub(r'{}.*?{}'.format(start, end), '', text, flags=re.DOTALL).strip()
    
    def SubElement(self, parent, tagname, **kwargs):
        """
        Helper for Paragraph bullet Point
        """
        element = OxmlElement(tagname)
        element.attrib.update(kwargs)
        parent.append(element)
        return element

    def getBulletInfo(self, paragraph, run=None):
        """
        Returns the attributes of the given <a:pPr> OxmlElement
        as well as its runs font-size.
        
        *param: paragraph* pptx _paragraph object
        *param: run* [optional] specific _run object
        """
        pPr = paragraph._p.get_or_add_pPr()
        if run is None:
            run = paragraph.runs[0]
        p_info = {
            "level": paragraph.level,
            "fontName": run.font.name,
            "fontSize": run.font.size,
        }
        return p_info

    def pBullet(
        self,
        paragraph,  # paragraph object
        font,  # fontName of that needs to be applied to bullet
        marL: int =864000,
        indent: int = -322920,
        size: int = 350000  # fontSize (in )
    ):
        """Bullets are set to Arial,
        actual text can be a different font
        """
        pPr = paragraph._p.get_or_add_pPr()
        # Set marL and indent attributes
        # Indent is the space between the bullet and the text.
        pPr.set('marL', str(marL))
        pPr.set('indent', str(indent))
        # Add buFont
        _ = self.SubElement(parent=pPr,
                    tagname="a:buSzPct",
                    val=str(size)
                    )
        _ = self.SubElement(parent=pPr,
                    tagname="a:buFont",
                    typeface=font,
                    )
        # Add buChar
        _ = self.SubElement(parent=pPr,
                    tagname='a:buChar',
                    char="•"
                    )

    def upload_file(self, file: UploadFile, user_id: str):
            """
            upload a file to the openwebui data base without the API (API doesn't work with the tools in version 0.6.5)
            ARGS:
                file: the file to upload
                user_id: the id of the user
            RETURNS:
                the file item
            """
            filename = file.filename
            id = str(uuid.uuid4())
            name = filename
            filename = f"{id}_{filename}"
            contents, file_path = Storage.upload_file(file.file, filename)

            file_item = Files.insert_new_file(
                user_id,
                FileForm(
                    **{
                        "id": id,
                        "filename": name,
                        "path": file_path,
                        "meta": {
                            "name": name,
                            "content_type": file.content_type,
                            "size": len(contents),
                            "data": {"generated_by": "upload_file"},
                        },
                    }
                ),
            )
            print(f"[DEBUG] File item: {file_item}")

            return file_item


    def add_title_slide(self, prs: Presentation, title: str ="Title", author: str = "author") -> None:
        """
        Creates and adds a new title slide to the given PowerPoint presentation.

        Args:
            prs (Presentation): The PowerPoint presentation object.
            title (str, optional): The title of the slide. Defaults to "Title".

        Returns:
            None

        Raises:
            ValueError: If the provided title is empty.
        """

        # Check if the title is not empty
        if not title.strip():
            raise ValueError("Slide title cannot be empty.")
        # Create new title slide
        slide_layout = prs.slide_layouts[self.slide_layouts["title_and_content"]]
        slide = prs.slides.add_slide(slide_layout)

        # fill in the content
        slide.shapes[0].text = title
        # add the author
        slide.shapes[1].text = author
        # add the date
        slide.shapes[3].text = datetime.now().strftime("%d/%m/%Y")
        


    def add_chapter_slide(
        self,
        prs: Presentation, 
        chapter: str = "Title", 
        suptitle: Optional[str] = None
    ) -> None:
        """
        Adds a new chapter title slide to a PowerPoint presentation.

        Args:
            prs (Presentation): The presentation object.
            chapter (str): The main title of the slide. Defaults to "Title".
            suptitle (str, optional): The suptitle of the slide. Defaults to None.

        Raises:
            ValueError: If the title is empty.

        Returns:
            None
        """
        # Check if the title is not empty
        if not chapter.strip():
            raise ValueError("Slide title cannot be empty.")
        # Create new title slide
        slide_layout = prs.slide_layouts[self.slide_layouts["chapter_title"]]
        slide = prs.slides.add_slide(slide_layout)

        # fill in the content
        slide.shapes[0].text = chapter
        if suptitle is not None and suptitle.strip():
            slide.shapes[1].text = suptitle

    def add_content_slide(
        self,
        prs: Presentation, 
        title: str = "Title", 
        content: str = "Content"
    ) -> None:
        """
        Adds a new title and content slide to the presentation.

        Args:
            prs (Presentation): The presentation object.
            title (str): The title of the slide. Defaults to "Title".
            content (str): The content of the slide. Defaults to "Content".

        Raises:
            ValueError: If the title or content is empty.
        """
        # Check if the title is not empty
        if not title.strip():
            raise ValueError("Slide title cannot be empty.")
        if not content.strip():
            raise ValueError("Slide content cannot be empty.")
        # Create new title slide
        slide_layout = prs.slide_layouts[self.slide_layouts["basic_content"]]
        slide = prs.slides.add_slide(slide_layout)

        # fill in the content
        slide.shapes[0].text = title
        body_shape = slide.shapes.placeholders[1]

        tf = body_shape.text_frame
        lines = content.split("\n")
        for line in lines:
            p = tf.add_paragraph()
            p.level = 0
            need_bullet = False
            spacing = 0
            while line.startswith('    '):
                line = line[4:]
                spacing += 1
            if line.startswith('* ') or line.startswith("• "):
                need_bullet = True
                line = line[2:]
            if need_bullet:
                p.text = line
                p.level = spacing + 1
                # For some reason, level 2 is dumb.
                # Should patch this in the template, but well ...
                if p.level >= 2:
                    p.level += 1
                
                print(self.getBulletInfo(p))
                #pBullet(p, "Arial", marL=p.level*864000, size=100000)
            else:
                p.text = '   '*spacing + line

    def add_final_slide(self, prs: Presentation, language: str) -> None:
        """
        Adds a final slide to the presentation.
        """
        if language == "fr":
            slide_layout = prs.slide_layouts[self.slide_layouts["final_slide_fr"]]
        else:
            slide_layout = prs.slide_layouts[self.slide_layouts["final_slide_en"]]
        slide = prs.slides.add_slide(slide_layout)
        


    
# --- Tools ---
class Tools:
    def __init__(self):
        """faire le truc avec le JSON """
        self.FILES_DIR = "./tmp"
        self.API_BASE_URL = "http://localhost:3000/api/v1/files/"
        # templates paths
        self.base_template_path = "./"
        self.fr_dir = "./"
        self.en_dir = "./"
        

        self.prefix= "CS-"
        os.makedirs(self.FILES_DIR, exist_ok=True)
        self.help_functions = HelpFunctions()
        self.event_emitter = EventEmitter()
    
    async def generate_pptx_from_json(self,language: str,confidentiality: str,json_data : dict,__request__: Request, __event_emitter__: Callable[[dict], Any] = None, __user__=None):
        """
        Generate a PowerPoint presentation from a JSON file.
        
        Args:
            language : The language of the presentation. (fr or en)
            confidentiality : The confidentiality of the presentation.  
            json_data : The JSON data to generate the presentation from.
            __user__ : The user to upload the file to.
        Returns:
            str: The download URL of the uploaded file.
        
        Example:
        Here is an example of the JSON data:
            
        ```json
            {
                "titre": "Titre de la présentation",
                "slides": [
                    {
                        "type": "titre",
                        "titre": "Titre de la présentation"
                    },
                    {
                        "type": "chapitre",
                        "titre": "Titre du chapitre"
                    },
                    {
                        "type": "contenu",
                        "titre": "Titre du chapitre",
                        "contenu": "Contenu de la slide\n* Bullet list\n.   * Bullet list niveau 2"
                    }
                ]
            }
        ```
        """
        emitter = EventEmitter(__event_emitter__)
        print("[DEBUG] json_data", json_data)
        # json_data = json.loads(json_data)
        # print("[DEBUG] json loaded", json_data)
        topic = json_data.get('titre')
        print("[DEBUG] topic", topic)
        await emitter.emit(f"Initiating pptx generation for topic: {topic}")

        # Load JSON data        
        print("[DEBUG] json_formatted", json_data)
        print("[DEBUG] language", language)
        print("[DEBUG] confidentiality", confidentiality)
        if confidentiality == "public":
            self.prefix = "CS-PU-"
        elif confidentiality == "internal":
            self.prefix = "CS-IN-"
        else:
            self.prefix = "CS-CO-"

        # Create presentation
        if language == "fr" or language == "french":
            # generating the template path french
            template_path = self.base_template_path + self.fr_dir + self.prefix + "template_fr.pptx"
            print("[DEBUG] french template path", template_path)
            prs = Presentation(template_path)
        else:
            # generating the template path english
            template_path = self.base_template_path + self.en_dir + self.prefix + "template_en.pptx"
            print("[DEBUG] english template path", template_path)
            prs = Presentation(template_path)


        print("[DEBUG] prs", prs)
        # Add title slide
        self.help_functions.add_title_slide(prs, title=json_data['titre'], author=__user__['name'])
        print("[DEBUG] title slide added")
        print("[DEBUG] user", __user__)

        # Add content slides
        try:
            await emitter.emit("Creating slides")
            for slide in json_data['slides']:
                print("[DEBUG] slide", slide)
                if slide['type'] == "chapitre":
                    self.help_functions.add_chapter_slide(prs, chapter=slide['titre'])
                    print("[DEBUG] chapter slide added")
                elif slide['type'] == "contenu":
                    self.help_functions.add_content_slide(prs, title=slide['titre'], content=slide['contenu'])
                    print("[DEBUG] content slide added")
            # add the final slide at the end of the presentation
            self.help_functions.add_final_slide(prs, language=language)

            await emitter.emit(
                status="complete",
                description=f"PPTX generation completed",
                done=True,
            )      
        except Exception as e:
            print("[DEBUG] Error", e)
            return "Error"
        
        
        # Save presentation
        if not os.path.exists(self.FILES_DIR):
            os.makedirs(self.FILES_DIR)
        # change the spaces to _ 
        json_data['titre'] = json_data['titre']
        # remove all special characters
        json_data['titre'] = re.sub(r'[^\w\s]', '', json_data['titre'])
        # remove all spaces
        json_data['titre'] = json_data['titre'].replace(' ', '_')

        output_path = self.FILES_DIR + '/' + self.prefix + json_data['titre'] + '.pptx'
        prs.save(output_path)
        print("[DEBUG] output_path", output_path)

        try :
            with open(output_path, 'rb') as f:
                print("[DEBUG] f", f)
                files = UploadFile(file=f, filename=os.path.basename(output_path))
                print("[DEBUG] files", files)
                file_item = await self.upload_file(file=files, user_id=__user__['id'] , __request__=__request__, __user__=__user__, __event_emitter__=__event_emitter__)
                print("[DEBUG] file_item", file_item)
                return file_item
        except Exception as e:
            print("[DEBUG] Error", e)
            return "Error"

    async def upload_file(self, file: UploadFile, user_id: str,__request__: Request,__user__:dict,__event_emitter__: Callable[[dict], Any] = None):

        emitter = EventEmitter(__event_emitter__)
        metadata = {"data": {"generated_by": "upload_file"}}
 
        await emitter.emit(f"getting download link for file : {file.filename}")
        
        # get the user for permissions
        user = Users.get_user_by_id(id=__user__['id'])
        print("[DEBUG] user", user)
        # upload the file to the database
        doc = upload_file(request=__request__, file=file, user=user, metadata=metadata, process=False)
        print("[DEBUG] doc", doc)

        # get the download link
        download_link = f"{self.API_BASE_URL}{doc.id}/content"
        print("[DEBUG] download_link", download_link)
        await emitter.emit(
                status="complete",
                description=f"finished generating the pptx file",
                done=True
            )
        return (
            f"<source><source_id>{doc.filename}</source_id><source_context>" 
            + str(download_link)
            + "</source_context></source>\n"
        )
    